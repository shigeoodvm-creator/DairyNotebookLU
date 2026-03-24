"""
子牛肺エコー検診パンフレット → PowerPoint生成スクリプト
A4縦（210mm × 297mm）1スライド
"""
from pptx import Presentation
from pptx.util import Mm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
import pptx.oxml.ns as nsmap
from lxml import etree

# ---------- ユーティリティ ----------
def rgb(r, g, b):
    return RGBColor(r, g, b)

def add_box(slide, left_mm, top_mm, w_mm, h_mm, fill_rgb=None, border_rgb=None, border_pt=0):
    """四角形を追加してShapeを返す"""
    from pptx.util import Mm
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Mm(left_mm), Mm(top_mm), Mm(w_mm), Mm(h_mm)
    )
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if border_rgb and border_pt > 0:
        line.color.rgb = border_rgb
        line.width = Pt(border_pt)
    else:
        line.fill.background()
    return shape

def add_text_box(slide, left_mm, top_mm, w_mm, h_mm, text, font_size=10,
                 bold=False, color=rgb(30,46,30), align=PP_ALIGN.LEFT,
                 wrap=True, font_name='Yu Gothic UI'):
    from pptx.util import Mm, Pt
    txb = slide.shapes.add_textbox(Mm(left_mm), Mm(top_mm), Mm(w_mm), Mm(h_mm))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txb

def add_multiline_textbox(slide, left_mm, top_mm, w_mm, h_mm, lines,
                           font_size=9, bold=False, color=rgb(30,46,30),
                           align=PP_ALIGN.LEFT, font_name='Yu Gothic UI',
                           line_spacing_pt=None):
    """複数行テキストボックス。linesはlist of (text, bold, size, color)"""
    from pptx.util import Mm, Pt
    from pptx.oxml.ns import qn
    txb = slide.shapes.add_textbox(Mm(left_mm), Mm(top_mm), Mm(w_mm), Mm(h_mm))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, b, sz, c = item, bold, font_size, color
        else:
            txt = item.get('text', '')
            b = item.get('bold', bold)
            sz = item.get('size', font_size)
            c = item.get('color', color)
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.bold = b
        run.font.color.rgb = c
        run.font.name = font_name
        if line_spacing_pt:
            from pptx.util import Pt as _Pt
            from pptx.oxml.ns import qn as _qn
            pPr = p._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
            spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
            spcPts.set('val', str(int(line_spacing_pt * 100)))
    return txb

# =============================================
# プレゼンテーション作成
# =============================================
prs = Presentation()

# A4縦に設定
prs.slide_width  = Mm(210)
prs.slide_height = Mm(297)

slide_layout = prs.slide_layouts[6]  # 空白レイアウト
slide = prs.slides.add_slide(slide_layout)

W = 210  # mm

# =============================================
# ① ヘッダー  top=0, h=70mm
#    緑グラデ背景 → 単色（PPTXはグラデ可だがシンプルに）
# =============================================
GREEN_DARK  = rgb(27, 92, 50)
GREEN_MID   = rgb(46, 128, 72)
WHITE       = rgb(255, 255, 255)
ACCENT_GRN  = rgb(160, 240, 192)
YELLOW      = rgb(255, 230, 100)
ORANGE      = rgb(255, 140, 0)
RED_SOFT    = rgb(220, 60, 60)
BLUE_SOFT   = rgb(50, 100, 200)
GRAY_LIGHT  = rgb(245, 248, 245)
BROWN_SOFT  = rgb(139, 90, 43)

# ヘッダー背景
hd = add_box(slide, 0, 0, W, 70, fill_rgb=GREEN_DARK)

# ヘッダー文字群（左側）
add_text_box(slide, 8, 5, 130, 12,
             '子牛 肺エコー検診サービス',
             font_size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_text_box(slide, 8, 19, 130, 8,
             '— 見えない肺炎を、早く・正確に発見 —',
             font_size=9, bold=False, color=ACCENT_GRN, align=PP_ALIGN.LEFT)

add_text_box(slide, 8, 30, 145, 20,
             '「うちの子牛、元気そうなのに\nなんで大きくならないんだろう…」\nそのお悩み、肺が原因かもしれません。',
             font_size=9.5, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

add_text_box(slide, 8, 52, 130, 10,
             '◯◯ 獣医師クリニック　子牛健康サポート',
             font_size=8, bold=False, color=ACCENT_GRN, align=PP_ALIGN.LEFT)

# ヘッダー：ちびキャラ子牛（SVGをPPTXで再現 → テキスト絵文字風 + 色付き楕円）
# PowerPointではSVGが使えないため、円・楕円・テキストで簡易キャラを作成
# 頭（大きい円）
COW_L = 155  # 子牛左端
# 体（小さい楕円）
body = add_box(slide, COW_L+8, 36, 30, 20, fill_rgb=rgb(245,245,245))
body.line.fill.background()
# 頭（大きい楕円）
head = add_box(slide, COW_L+3, 18, 38, 28, fill_rgb=WHITE)
head.line.fill.background()
# 耳（左）
ear_l = add_box(slide, COW_L+1, 19, 9, 7, fill_rgb=rgb(255,220,200))
ear_l.line.fill.background()
# 耳（右）
ear_r = add_box(slide, COW_L+34, 19, 9, 7, fill_rgb=rgb(255,220,200))
ear_r.line.fill.background()
# 鼻（楕円）
nose = add_box(slide, COW_L+12, 36, 18, 9, fill_rgb=rgb(255,200,200))
nose.line.fill.background()
# 前足
leg1 = add_box(slide, COW_L+10, 53, 7, 12, fill_rgb=rgb(240,240,240))
leg1.line.fill.background()
leg2 = add_box(slide, COW_L+25, 53, 7, 12, fill_rgb=rgb(240,240,240))
leg2.line.fill.background()
# 模様
spot = add_box(slide, COW_L+14, 22, 12, 10, fill_rgb=rgb(220,220,220))
spot.line.fill.background()
spot2 = add_box(slide, COW_L+12, 38, 8, 8, fill_rgb=rgb(220,220,220))
spot2.line.fill.background()
# 目（テキストで）
add_text_box(slide, COW_L+8, 23, 8, 6, '●', font_size=9, bold=True, color=rgb(80,40,20))
add_text_box(slide, COW_L+27, 23, 8, 6, '●', font_size=9, bold=True, color=rgb(80,40,20))
# ハート
add_text_box(slide, COW_L+17, 13, 12, 7, '♡', font_size=10, bold=True, color=rgb(255,100,150))
# キラキラ
add_text_box(slide, COW_L+36, 11, 10, 7, '✦', font_size=7, color=YELLOW)
add_text_box(slide, COW_L+0, 12, 10, 7, '✦', font_size=6, color=YELLOW)

# =============================================
# ② 問題提起  top=70, h=38mm
# =============================================
add_box(slide, 0, 70, W, 38, fill_rgb=rgb(255,252,240))

add_text_box(slide, 8, 73, W-16, 9,
             '⚠ 「不顕性肺炎」を知っていますか？',
             font_size=12, bold=True, color=rgb(180,80,0), align=PP_ALIGN.LEFT)

add_multiline_textbox(slide, 8, 83, W-16, 22,
    [
        {'text': '子牛の肺炎は「咳・鼻水・発熱」が出て初めて気づくケースがほとんどです。', 'size': 9},
        {'text': 'しかし実は、症状が出る前から肺の中で炎症が進行していることが多くあります。', 'size': 9},
        {'text': '肺エコーを使えば、こうした"見えない肺炎"を症状が出る前に発見できます。', 'size': 9, 'bold': True, 'color': rgb(180,80,0)},
    ],
    line_spacing_pt=14)

# =============================================
# ③ 経済的損失  top=108, h=50mm
# =============================================
add_box(slide, 0, 108, W, 50, fill_rgb=WHITE)

add_text_box(slide, 8, 111, W-16, 9,
             '📉 放置すると…こんなリスクがあります',
             font_size=11, bold=True, color=RED_SOFT, align=PP_ALIGN.LEFT)

# 3列カード
card_data = [
    ('💊', '治療費増大', '発見が遅れるほど\n治療が長引き\nコストが増加'),
    ('📏', '発育遅延', '肺炎既往牛は\n発育が遅く\n出荷体重が落ちやすい'),
    ('⚡', '廃用リスク', '重症化すると\n慢性肺炎となり\n早期廃用の原因に'),
]
for i, (icon, title, body_txt) in enumerate(card_data):
    cx = 8 + i * 64
    add_box(slide, cx, 121, 60, 33, fill_rgb=rgb(255,245,245), border_rgb=rgb(220,180,180), border_pt=0.5)
    add_text_box(slide, cx+2, 123, 56, 9, f'{icon} {title}', font_size=9.5, bold=True, color=RED_SOFT)
    add_text_box(slide, cx+2, 132, 56, 20, body_txt, font_size=8.5, color=rgb(60,40,40))

# =============================================
# ④ サービスの流れ  top=158, h=48mm
# =============================================
add_box(slide, 0, 158, W, 48, fill_rgb=rgb(240,248,240))

add_text_box(slide, 8, 161, W-16, 9,
             '🐄 検診の流れ',
             font_size=11, bold=True, color=GREEN_DARK, align=PP_ALIGN.LEFT)

steps = [
    ('①', '事前\nご連絡'),
    ('②', '農場\n訪問'),
    ('③', '超音波\n検査'),
    ('④', 'その場で\n結果説明'),
    ('⑤', 'データ\n報告書'),
]
sw = (W - 16) / 5  # 各ステップ幅
for i, (num, label) in enumerate(steps):
    sx = 8 + i * sw
    add_box(slide, sx+1, 172, sw-4, 28, fill_rgb=GREEN_MID)
    add_text_box(slide, sx+1, 173, sw-4, 8, num, font_size=9, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    add_text_box(slide, sx+1, 181, sw-4, 16, label, font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 矢印（最後以外）
    if i < 4:
        add_text_box(slide, sx+sw-3, 180, 6, 8, '▶', font_size=7, color=GREEN_DARK)

add_text_box(slide, 8, 201, W-16, 7,
             '※ 検査は農場に伺い、子牛1頭ずつ超音波プローブを当てるだけ。拘束・痛みはありません。',
             font_size=7.5, color=rgb(60,100,60))

# =============================================
# ⑤ 料金  top=206, h=52mm
# =============================================
add_box(slide, 0, 206, W, 52, fill_rgb=WHITE)

add_text_box(slide, 8, 209, W-16, 9,
             '💴 料金のご案内',
             font_size=11, bold=True, color=GREEN_DARK, align=PP_ALIGN.LEFT)

# 料金表
fee_data = [
    ('農場訪問料', '11,250 円', '（交通費・出張費込み）'),
    ('データ分析料', '7,500 円', '（個体別レポート作成）'),
]
for i, (item, price, note) in enumerate(fee_data):
    fy = 220 + i * 16
    add_box(slide, 8, fy, W-16, 14, fill_rgb=rgb(240,250,240), border_rgb=rgb(100,160,100), border_pt=0.5)
    add_text_box(slide, 12, fy+2, 70, 10, item, font_size=10, bold=True, color=GREEN_DARK)
    add_text_box(slide, 85, fy+1, 55, 12, price, font_size=13, bold=True, color=rgb(180,80,0), align=PP_ALIGN.RIGHT)
    add_text_box(slide, 143, fy+3, 60, 8, note, font_size=7.5, color=rgb(80,80,80))

# 合計
add_box(slide, 8, 252, W-16, 3, fill_rgb=GREEN_DARK)
add_text_box(slide, 8, 255, W-16, 5, '', font_size=1, color=WHITE)  # spacer

add_box(slide, 8, 253, W-16, 0, fill_rgb=WHITE)
add_text_box(slide, 12, 254, 80, 8, '合計', font_size=11, bold=True, color=GREEN_DARK)
add_text_box(slide, 90, 253, 80, 10, '18,750 円', font_size=16, bold=True, color=rgb(180,80,0), align=PP_ALIGN.RIGHT)
add_text_box(slide, 173, 256, 30, 7, '（税込）', font_size=8, color=rgb(80,80,80))

# =============================================
# ⑥ フッター  top=258, h=39mm
# =============================================
add_box(slide, 0, 258, W, 39, fill_rgb=GREEN_DARK)

add_text_box(slide, 0, 261, W, 9,
             'お問い合わせ・ご予約',
             font_size=10, bold=True, color=ACCENT_GRN, align=PP_ALIGN.CENTER)

add_text_box(slide, 0, 271, W, 9,
             '◯◯ 獣医師',
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, 0, 280, W, 8,
             'TEL: 000-0000-0000',
             font_size=11, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

add_text_box(slide, 0, 289, W, 7,
             '※ 連絡先・クリニック名を上記に記入してご使用ください',
             font_size=7, color=ACCENT_GRN, align=PP_ALIGN.CENTER)

# =============================================
# 保存
# =============================================
out_path = r'C:\Users\user\OneDrive\デスクトップ\肺エコープロジェクト\子牛肺エコー検診_パンフレット.pptx'
prs.save(out_path)
print(f'保存完了: {out_path}')
